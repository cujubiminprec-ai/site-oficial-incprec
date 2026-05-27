#!/usr/bin/env python3
"""
Convert PPTX slides to JPEG images using python-pptx + Pillow.
Usage: python3 pptx_to_slides.py <pptx_path> <out_dir> <prefix>
Output: JSON with slide_count and array of saved image paths (/uploads/slide-previews/...).
"""
import sys, json, os, io

try:
    from pptx import Presentation
    from pptx.util import Pt
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_OK = True
except ImportError:
    PIL_OK = False


def get_bg_color(slide):
    try:
        fill = slide.background.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and hasattr(fc, "rgb"):
                c = fc.rgb
                return (c.red, c.green, c.blue)
    except Exception:
        pass
    return (22, 78, 43)  # INPREC green fallback


def load_font(size, bold=False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf" if bold else "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                pass
    return ImageFont.load_default()


def wrap_text(text, font, max_width, draw):
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = (current + " " + word).strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_text_placeholder(img, slide, bg):
    draw = ImageDraw.Draw(img)
    W, H = img.size

    r, g, b = bg
    lum = 0.299 * r + 0.587 * g + 0.114 * b
    text_color = (255, 255, 255) if lum < 128 else (20, 20, 20)
    muted_color = (200, 200, 200) if lum < 128 else (80, 80, 80)

    # Draw subtle grid overlay for visual interest
    grid_color = (r + 20, g + 20, b + 20) if lum < 128 else (r - 15, g - 15, b - 15)
    grid_color = tuple(max(0, min(255, c)) for c in grid_color)
    for x in range(0, W, 40):
        draw.line([(x, 0), (x, H)], fill=grid_color, width=1)
    for y in range(0, H, 40):
        draw.line([(0, y), (W, y)], fill=grid_color, width=1)

    # Collect text from shapes
    texts = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        except Exception:
            pass

    if not texts:
        # No text: just show slide number indicator
        font_lg = load_font(48, bold=True)
        label = "SLIDE"
        bbox = draw.textbbox((0, 0), label, font=font_lg)
        tw = bbox[2] - bbox[0]
        draw.text(((W - tw) // 2, H // 2 - 30), label, fill=muted_color, font=font_lg)
        return

    font_title = load_font(38, bold=True)
    font_body = load_font(22, bold=False)

    padding = 60
    max_w = W - padding * 2
    y = padding

    # Title (first text block)
    title_lines = wrap_text(texts[0][:120], font_title, max_w, draw)
    for line in title_lines[:3]:
        bbox = draw.textbbox((0, 0), line, font=font_title)
        th = bbox[3] - bbox[1]
        draw.text((padding, y), line, fill=text_color, font=font_title)
        y += th + 8

    # Separator
    if y < H - 100:
        y += 14
        draw.line([(padding, y), (W - padding, y)], fill=muted_color, width=2)
        y += 18

    # Body texts
    for body_text in texts[1:6]:
        if y > H - 60:
            break
        body_lines = wrap_text(body_text[:200], font_body, max_w, draw)
        for line in body_lines[:2]:
            if y > H - 60:
                break
            bbox = draw.textbbox((0, 0), line, font=font_body)
            th = bbox[3] - bbox[1]
            draw.text((padding, y), line, fill=muted_color, font=font_body)
            y += th + 6
        y += 4


def extract_slide_image(slide, idx, out_dir, prefix, slide_w_emu, slide_h_emu):
    TARGET_W = 960
    TARGET_H = 540

    bg = get_bg_color(slide)
    canvas = Image.new("RGB", (TARGET_W, TARGET_H), bg)

    # Find all picture shapes, pick the one with the largest area
    pictures = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                area = shape.width * shape.height
                pictures.append((area, shape))
        except Exception:
            pass
    pictures.sort(key=lambda x: x[0], reverse=True)

    if pictures:
        _, best = pictures[0]
        try:
            img_bytes = best.image.blob
            src = Image.open(io.BytesIO(img_bytes)).convert("RGB")

            # Scale to fill the canvas
            src_ratio = src.width / src.height
            canvas_ratio = TARGET_W / TARGET_H

            if src_ratio > canvas_ratio:
                new_w = TARGET_W
                new_h = int(TARGET_W / src_ratio)
            else:
                new_h = TARGET_H
                new_w = int(TARGET_H * src_ratio)

            src = src.resize((new_w, new_h), Image.LANCZOS)
            paste_x = (TARGET_W - new_w) // 2
            paste_y = (TARGET_H - new_h) // 2
            canvas.paste(src, (paste_x, paste_y))
        except Exception:
            draw_text_placeholder(canvas, slide, bg)
    else:
        draw_text_placeholder(canvas, slide, bg)

    filename = f"{prefix}_slide_{idx + 1:03d}.jpg"
    path = os.path.join(out_dir, filename)
    canvas.save(path, "JPEG", quality=88, optimize=True)
    return filename


def convert(pptx_path, out_dir, prefix):
    if not PPTX_OK:
        return {"success": False, "error": "python-pptx not installed"}
    if not PIL_OK:
        return {"success": False, "error": "Pillow not installed"}

    if not os.path.exists(pptx_path):
        return {"success": False, "error": f"File not found: {pptx_path}"}

    os.makedirs(out_dir, exist_ok=True)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"success": False, "error": f"Cannot open PPTX: {str(e)}"}

    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 5143500

    images = []
    errors = []
    for i, slide in enumerate(prs.slides):
        try:
            fname = extract_slide_image(slide, i, out_dir, prefix, slide_w, slide_h)
            images.append(f"/uploads/slide-previews/{fname}")
        except Exception as e:
            errors.append(f"slide {i+1}: {str(e)}")

    return {
        "success": len(images) > 0,
        "slide_count": len(prs.slides),
        "converted": len(images),
        "slide_images": images,
        "errors": errors,
    }


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print(json.dumps({"success": False, "error": "Usage: pptx_to_slides.py <pptx_path> <out_dir> <prefix>"}))
        sys.exit(1)
    result = convert(sys.argv[1], sys.argv[2], sys.argv[3])
    print(json.dumps(result))
    sys.exit(0 if result.get("success") else 1)
